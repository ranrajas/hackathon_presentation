from pptx import Presentation
from pptx.util import Inches, Pt
import uuid

def create_ppt(input_data):
    print (input_data)

    prs = Presentation()

    # --- Title Slide ---
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = input_data["title"]

    # --- Content Slides ---
    for slide_data in input_data["slides"]:
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)

        slide.shapes.title.text = slide_data["title"]

        body = slide.shapes.placeholders[1]
        tf = body.text_frame
        tf.clear()

        for point in slide_data["points"]:
            p = tf.add_paragraph()
            p.text = point
            p.font.size = Pt(18)
            p.level = 1

    # --- Save file ---
    filename = f"presentation_{uuid.uuid4().hex}.pptx"
    prs.save(filename)

    return filename
